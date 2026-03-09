import streamlit as st
import pandas as pd
import akasha
import os
import akasha.helper as ah
import shutil
import json
import docx2txt
from pptx import Presentation
from pypdf import PdfReader
import sys
import re
from datetime import datetime
import traceback


st.set_page_config(page_title="CSW")
# 自訂樣式：讓 Markdown 的程式碼區塊自動換行、避免左右滑動
st.markdown(
        """
        <style>
        div[data-testid="stMarkdownContainer"] pre {
            white-space: pre-wrap !important;
            word-wrap: break-word !important;
            word-break: break-word !important;
            overflow-x: hidden !important;
        }
        div[data-testid="stMarkdownContainer"] code {
            white-space: pre-wrap !important;
            word-wrap: break-word !important;
            word-break: break-word !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
)
# --- 1. 環境設定 ---
DATA_FOLDER = os.getenv("DATA_FOLDER", "data")
DEFAULT_DATA_FILE = os.getenv("DEFAULT_DATA_FILE", "default_data/FAQ_Default.xlsx")
os.makedirs(DATA_FOLDER, exist_ok=True)
DEFAULT_FILE = os.path.join(DATA_FOLDER, "FAQ_Default.xlsx")
DATA_STATE_PATH = "data_state.json"
CHAT_LOGS_FOLDER = os.path.join(DATA_FOLDER, "chat_logs")
os.makedirs(CHAT_LOGS_FOLDER, exist_ok=True)
ALLOWED_EXTS = {".xlsx", ".docx", ".txt", ".pdf", ".pptx"}
if not os.path.exists(DEFAULT_FILE):
    if not os.path.exists(DEFAULT_DATA_FILE):
        st.write(f"缺少預設文件{DEFAULT_DATA_FILE}，請建立資料夾 default_data 並將 FAQ_Default.xlsx 存入後重新整理頁面。")
        st.stop()
    else:
        shutil.copy(DEFAULT_DATA_FILE, DEFAULT_FILE)

MODEL_CONFIG = {
    "Google Gemini(2.5-flash)": {
        "env_var": "GEMINI_API_KEY",
        "model_name": "gemini:gemini-2.5-flash"
    },
    "OpenAI (GPT-4o)": {
        "env_var": "OPENAI_API_KEY",
        "model_name": "openai:gpt-4o"
    },
    "OpenAI (GPT-5)": {
        "env_var": "OPENAI_API_KEY",
        "model_name": "openai:gpt-5"
    },
    "Anthropic Claude": {
        "env_var": "ANTHROPIC_API_KEY",
        "model_name": "claude:claude-3-opus-20240229"
    }
}

# 初始化 Session State
if "history_list" not in st.session_state:
    st.session_state.history_list = []
if "messages" not in st.session_state:
    st.session_state.messages = []
# 透過 data_state.json 管理目前是否使用預設檔與檔名
def load_data_state():
    try:
        if os.path.exists(DATA_STATE_PATH):
            with open(DATA_STATE_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {"mode": "default", "file_name": ["FAQ_Default.xlsx"]}

def save_data_state(mode: str, file_names):
    try:
        files = file_names if isinstance(file_names, list) else [file_names]
        state = {}
        if os.path.exists(DATA_STATE_PATH):
            with open(DATA_STATE_PATH, "r", encoding="utf-8") as f:
                state = json.load(f) or {}
        state["mode"] = mode
        state["file_name"] = files
        with open(DATA_STATE_PATH, "w", encoding="utf-8") as f:
            json.dump(state, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def get_chat_active_file():
    try:
        state = load_data_state() or {}
        chat = state.get("chat_state") or {}
        return chat.get("active_file")
    except Exception:
        return None

def set_chat_active_file(path: str):
    try:
        state = {}
        if os.path.exists(DATA_STATE_PATH):
            with open(DATA_STATE_PATH, "r", encoding="utf-8") as f:
                state = json.load(f) or {}
        state["chat_state"] = {"active_file": path}
        with open(DATA_STATE_PATH, "w", encoding="utf-8") as f:
            json.dump(state, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

state_json = load_data_state()
if "include_default" not in st.session_state:
    st.session_state.include_default = (state_json.get("mode") == "default")
if "use_data_name" not in st.session_state:
    files = state_json.get("file_name") or state_json.get("file") or []
    st.session_state.use_data_name = files if files else []
if "current_data" not in st.session_state:
    st.session_state.current_data = None
if "processed_files" not in st.session_state:
    st.session_state.processed_files = []

# 假設圖片路徑
AVATAR_PATH = "static"
CSW_AVATAR = os.path.join(AVATAR_PATH, "csw_icon.jpg")
USER_AVATAR = os.path.join(AVATAR_PATH, "user_icon.jpg")
# BOT_AVATAR = "https://your-domain.com/bot-logo.png"
# with st.chat_message("user", avatar=USER_AVATAR):

# --- 2. 工具函數 ---
@st.cache_data    
def read_excel_sheets(file_path):
    if not os.path.exists(file_path):
        return None
    try:
        # 自動讀取所有工作表
        return pd.read_excel(file_path, sheet_name=None)
    except Exception as e:
        st.error(f"讀取 Excel 失敗: {e}")
        return None

@st.cache_data
def extract_text_from_docx(path: str) -> str:
    try:
        return docx2txt.process(path) or ""
    except Exception as e:
        st.error(f"讀取 DOCX 失敗: {e}")
        return ""

@st.cache_data
def extract_text_from_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        st.error(f"讀取 TXT 失敗: {e}")
        return ""

@st.cache_data
def extract_text_from_pdf(path: str) -> str:
    try:
        reader = PdfReader(path)
        pages_text = []
        for p in reader.pages:
            try:
                pages_text.append(p.extract_text() or "")
            except Exception:
                pages_text.append("")
        return "\n".join(pages_text)
    except Exception as e:
        st.error(f"讀取 PDF 失敗: {e}")
        return ""

@st.cache_data
def extract_text_from_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    texts.append("\n".join([p.text for p in shape.text_frame.paragraphs]))
        return "\n".join(texts)
    except Exception as e:
        st.error(f"讀取 PPTX 失敗: {e}")
        return ""

def df_from_text(text: str, source_label: str) -> pd.DataFrame:
    return pd.DataFrame({"來源": [source_label], "內容": [text]})

def read_excel_list(file_paths):
    if not file_paths:
        return None
    combined = {}
    for p in file_paths:
        ext = os.path.splitext(p)[1].lower()
        base = os.path.basename(p)
        if ext == ".xlsx":
            data = read_excel_sheets(p)
            if not data:
                continue
            for sheet_name, df in data.items():
                # 保留原始檔案名稱 + 工作表名稱，避免不同檔案同名工作表混合
                key = f"{base}::{sheet_name}"
                if key in combined:
                    try:
                        combined[key] = pd.concat([combined[key], df], ignore_index=True)
                    except Exception as e:
                        st.error(f"合併工作表 {sheet_name} 失敗: {e}")
                else:
                    combined[key] = df
        elif ext == ".docx":
            text = extract_text_from_docx(p)
            df = df_from_text(text, base)
            key = base
            combined[key] = pd.concat([combined[key], df], ignore_index=True) if key in combined else df
        elif ext == ".txt":
            text = extract_text_from_txt(p)
            df = df_from_text(text, base)
            key = base
            combined[key] = pd.concat([combined[key], df], ignore_index=True) if key in combined else df
        elif ext == ".pdf":
            text = extract_text_from_pdf(p)
            df = df_from_text(text, base)
            key = base
            combined[key] = pd.concat([combined[key], df], ignore_index=True) if key in combined else df
        elif ext == ".pptx":
            text = extract_text_from_pptx(p)
            df = df_from_text(text, base)
            key = base
            combined[key] = pd.concat([combined[key], df], ignore_index=True) if key in combined else df
        else:
            st.warning(f"不支援的檔案類型: {base}")
    return combined if combined else None

def format_data_for_ai(data_dict):
    """將 DataFrame 字典轉為 AI 易讀的字串"""
    if not data_dict: return "目前無參考資料。"
    full_text = ""
    for name, df in data_dict.items():
        full_text += f"\n--- {name} ---\n"
        full_text += df.to_csv(index=False)
    return full_text

def normalize_response_text(text):
    """將回應中的可見換行符（如 \\n\\n、\\r\\n）轉為真正的換行，並壓縮過多的空行。"""
    if text is None:
        return ""
    if not isinstance(text, str):
        text = str(text)
    # 標準化行結尾
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # 將字面上的反斜線換行轉為真正的換行
    text = text.replace("\\r\\n", "\n").replace("\\n", "\n")
    # 壓縮連續 3 行以上空行為 2 行
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text

def sanitize_filename(name: str) -> str:
    """移除檔名中的不合法字元，保留英數字、中文、空白、-與_，並壓縮重複空白。"""
    if not name:
        return ""
    # 去除換行與引號等不必要符號
    name = re.sub(r"[`\"'\r\n]+", " ", name).strip()
    # 只允許中英文、數字、空白、-、_
    name = re.sub(r"[^\w\u4e00-\u9fff\s\-]", " ", name)
    # 壓縮空白
    name = re.sub(r"\s+", " ", name).strip()
    # 避免空字串
    return name[:80] if name else ""

def generate_unique_log_path(base_name: str) -> str:
    """根據檔名產生唯一的 JSON 檔案路徑（於 CHAT_LOGS_FOLDER 下）。"""
    base_name = base_name.rstrip(".json")
    candidate = f"{base_name}.json"
    path = os.path.join(CHAT_LOGS_FOLDER, candidate)
    if not os.path.exists(path):
        return path
    idx = 2
    while True:
        candidate = f"{base_name}-{idx}.json"
        path = os.path.join(CHAT_LOGS_FOLDER, candidate)
        if not os.path.exists(path):
            return path
        idx += 1

def generate_chat_filename_path(hint: str | None = None) -> str:
    """使用語言模型產生對話記錄檔名，並回傳唯一的完整路徑。"""
    try:
        model_name = config["model_name"]
        ask_fn = akasha.ask(
            model=model_name,
            temperature=1.0,
            max_input_tokens=2000,
            max_output_tokens=2000,
        )
        prompt = (
            """
            <任務>\n
            請根據使用者向系統的第一句提問內容，產生一個適合顯示給一般使用者看的對話紀錄名稱。\n
            </任務>\n
            <規則>\n
            1. 名稱長度限制在 6 至 16 個中文字以內或 30 個英文字以內\n
            2. 不要出現標點符號、引號或特殊符號\n
            3. 不要包含日期、時間、編號\n
            4. 以「問題主題或使用者意圖」作為命名重點\n
            5. 避免照抄原句，請進行語意摘要或重述\n
            </規則>\n
        """
        )
        if hint:
            prompt += f"- 使用者的第一句提問內容：{hint}\n"
        name_raw = ask_fn(prompt=prompt)
        name = normalize_response_text(name_raw).strip().splitlines()[0]
        safe = sanitize_filename(name)
        if not safe:
            safe = "對話紀錄"
        return generate_unique_log_path(safe)
    except Exception:
        # LM 失敗時的回退方案
        return generate_unique_log_path("對話紀錄")

def save_chat_log(create_if_missing: bool = True):
    try:
        path = get_chat_active_file()
        started_at = None
        if not path and create_if_missing:
            # 以語言模型產生檔名，若可用則參考最後一則使用者訊息
            last_user = None
            try:
                msgs = st.session_state.get("messages", []) or []
                for m in reversed(msgs):
                    if m.get("role") == "user":
                        last_user = m.get("content")
                        break
            except Exception:
                pass
            path = generate_chat_filename_path(last_user)
            set_chat_active_file(path)
            started_at = datetime.now().isoformat()

        if not path:
            return None

        data = {
            "timestamp": datetime.now().isoformat(),
            "messages": st.session_state.get("messages", []),
            "history_list": st.session_state.get("history_list", []),
        }
        # 保留或設定對話開始時間
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    prev = json.load(f)
                if "started_at" in prev:
                    data["started_at"] = prev["started_at"]
            except Exception:
                pass
        if started_at and "started_at" not in data:
            data["started_at"] = started_at

        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return path
    except Exception:
        return None

def list_chat_logs():
    try:
        files = [fn for fn in os.listdir(CHAT_LOGS_FOLDER) if fn.lower().endswith(".json")]
        return sorted(files, reverse=True)
    except Exception:
        return []

def load_chat_log(file_name: str):
    try:
        path = os.path.join(CHAT_LOGS_FOLDER, file_name)
        if not os.path.exists(path):
            return None
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None

def start_new_conversation():
    """開始新對話但不立即建立檔案；於首次訊息回覆時才建立記錄。"""
    # 清空目前的訊息與歷史
    st.session_state.messages = []
    st.session_state.history_list = []
    # 清除目前的對話檔案指標
    set_chat_active_file("")
    return None

# 定義一個內部函數來把 list 轉回字串，方便計算 Token
def get_history_string(h_list):
    return "".join([f"\n提問: {item['q']}\n回覆: {item['a']}" for item in h_list])

# --- 3. 初始資料載入邏輯 ---
# 只有在 current_data 是 None 的時候才去執行讀取，並依照 toggle 狀態決定來源
if st.session_state.current_data is None:
    # 依照 checkbox 選擇決定載入資料
    files = st.session_state.use_data_name if isinstance(st.session_state.use_data_name, list) else []
    paths = [os.path.join(DATA_FOLDER, f) for f in files if os.path.exists(os.path.join(DATA_FOLDER, f))]
    if st.session_state.include_default:
        paths = [DEFAULT_FILE] + paths
    if paths:
        st.session_state.current_data = read_excel_list(paths)
        save_data_state("default" if (st.session_state.include_default and not files) else "active", files)
    else:
        # 無選擇時使用預設
        st.session_state.include_default = True
        st.session_state.current_data = read_excel_sheets(DEFAULT_FILE)
        save_data_state("default", ["FAQ_Default.xlsx"])

# --- 4. Streamlit 側邊欄介面設定 ---
with st.sidebar:
    with st.expander("模型與 API 設定", expanded=True):
        # 1.下拉式選單選擇模型
        selected_model_display = st.selectbox("選擇模型來源",options=list(MODEL_CONFIG.keys()))
        # 取得對應的配置
        config = MODEL_CONFIG[selected_model_display]

        # 2.加入API_KEY輸入框
        user_api_key = st.text_input(
            "輸入您的 API KEY", 
            type="password",
            help="輸入有效API_KEY後即可進行對話"
        )
        api_valid = False
        if user_api_key:
            os.environ[config["env_var"]] = user_api_key
            # 發送一次測試請求以確認 Key 有效性
            try:
                test_ak = akasha.ask(
                    model=config["model_name"],
                    temperature=0.1,
                )
                test = test_ak(prompt="return hi")
                st.success("API Key 已就緒！")
                api_valid = True 
            except Exception as e:
                st.error(f"API Key 無效，請檢查後重新輸入。")
                api_valid = False
        else:
            st.warning("請先輸入 API Key")
        
    # 對話組管理
    with st.expander("對話組", expanded=False):
        current_file = get_chat_active_file()
        current_name_raw = os.path.basename(current_file) if current_file else "尚未選擇"
        current_name = (
            current_name_raw[:-5]
            if isinstance(current_name_raw, str) and current_name_raw.lower().endswith(".json")
            else current_name_raw
        )
        st.caption(f"目前對話檔案：{current_name}")

        if st.button("開啟新對話", key="btn_new_conversation"):
            start_new_conversation()
            st.rerun()

        logs_real = list_chat_logs()
        NONE_OPTION = "__NONE__"
        options = [NONE_OPTION] + logs_real

        # 讓選擇框同步目前活躍對話；若無則選擇占位項
        active_file = get_chat_active_file()
        active_name = os.path.basename(active_file) if active_file else None
        desired_selection = active_name if (active_name and active_name in logs_real) else NONE_OPTION
        try:
            cur_sel = st.session_state.get("sel_chat_group")
            if (not cur_sel) or (cur_sel not in options) or (cur_sel == NONE_OPTION and desired_selection != NONE_OPTION):
                st.session_state["sel_chat_group"] = desired_selection
        except Exception:
            pass

        # 當使用者變更選擇時才載入（避免初始渲染就載入）
        def _on_select_chat():
            try:
                sel = st.session_state.get("sel_chat_group")
                if sel == NONE_OPTION:
                    # 清除活躍檔案指標，但不動現有畫面訊息
                    set_chat_active_file("")
                    return
                cur = os.path.basename(get_chat_active_file() or "") or None
                if sel and sel != cur:
                    data = load_chat_log(sel)
                    if data:
                        st.session_state.messages = data.get("messages", [])
                        st.session_state.history_list = data.get("history_list", [])
                        set_chat_active_file(os.path.join(CHAT_LOGS_FOLDER, sel))
            except Exception:
                pass

        selected_log = st.selectbox(
            "選擇對話載入",
            options=options,
            format_func=lambda s: ("（未選擇）" if s == NONE_OPTION else (s[:-5] if s.lower().endswith(".json") else s)),
            key="sel_chat_group",
            on_change=_on_select_chat,
        )

        if selected_log and selected_log != NONE_OPTION:
            preview = load_chat_log(selected_log) or {}
            msg_count = len(preview.get("messages") or [])
            name_display = selected_log[:-5] if selected_log.lower().endswith(".json") else selected_log
            st.caption(f"訊息數：{msg_count} | 檔名：{name_display}")
            # 更名（上）
            new_name_default = name_display
            new_name = st.text_input(
                "重新命名對話",
                value=new_name_default,
                key=f"rename_input_{selected_log}"
            )
            if st.button("確認重新命名對話", key=f"btn_rename_{selected_log}"):
                try:
                    base = sanitize_filename(new_name)
                    if not base:
                        st.warning("請輸入有效的名稱（僅中英文、數字、空白、-或_）。")
                    else:
                        src = os.path.join(CHAT_LOGS_FOLDER, selected_log)
                        dst = generate_unique_log_path(base)
                        if src == dst:
                            st.info("名稱未變更。")
                        else:
                            os.rename(src, dst)
                            # 若目前為活躍對話則更新指向
                            active_file = get_chat_active_file()
                            if active_file and os.path.basename(active_file) == selected_log:
                                set_chat_active_file(dst)
                            st.success(f"✅ 已更名為：{os.path.basename(dst)[:-5]}")
                            st.rerun()
                except Exception as e:
                    st.error(f"更名失敗：{e}")

            # 刪除（下）
            st.caption("刪除目前對話紀錄")
            if st.button("刪除目前對話紀錄", key=f"del_{selected_log}", help="注意：刪除後無法復原！", type="primary"):
                try:
                    target = os.path.join(CHAT_LOGS_FOLDER, selected_log)
                    if os.path.exists(target):
                        os.remove(target)
                        # 若刪除的是當前對話檔案，重置當前對話狀態
                        current_file = get_chat_active_file()
                        if current_file and os.path.basename(current_file) == selected_log:
                            set_chat_active_file("")
                            st.session_state.messages = []
                            st.session_state.history_list = []
                        st.success(f"🗑️ 已刪除對話：{selected_log}")
                        st.rerun()
                    else:
                        st.warning("找不到檔案，可能已被刪除或移動。")
                except Exception as e:
                    st.error(f"刪除失敗：{e}")
        else:
            st.caption("尚未選擇對話。您可以從下拉選單選擇或開啟新對話。")

    # 3.資料上傳（合併至摺疊區塊內）

    # 使用預設資料庫選項與檔案勾選（摺疊區塊）
    # 列出可用檔案並提供勾選
    def list_available_files():
        try:
            files = []
            for fn in os.listdir(DATA_FOLDER):
                path = os.path.join(DATA_FOLDER, fn)
                if os.path.isfile(path) and os.path.splitext(fn)[1].lower() in ALLOWED_EXTS and fn != os.path.basename(DEFAULT_FILE):
                    files.append(fn)
            return sorted(files)
        except Exception:
            return []

    with st.expander("選擇生效檔案與刪除", expanded=False):
        # 上傳更新資料
        uploaded_files = st.file_uploader(
            "上傳更新資料 (xlsx/docx/txt/pdf/pptx)", 
            type=["xlsx", "docx", "txt", "pdf", "pptx"],
            accept_multiple_files=True,
            )
        if uploaded_files:
            # 只處理尚未儲存過的新檔案（以檔名判斷）
            processed = st.session_state.processed_files if isinstance(st.session_state.processed_files, list) else []
            new_uploads = [uf for uf in uploaded_files if uf.name not in processed]
            if new_uploads:
                saved_names = []
                for uf in new_uploads:
                    target_path = os.path.join(DATA_FOLDER, uf.name)
                    with open(target_path, "wb") as f:
                        f.write(uf.getbuffer())
                    saved_names.append(uf.name)
                # 更新已處理名單
                st.session_state.processed_files = list(dict.fromkeys(processed + saved_names))
                # 更新目前的使用清單：保留原有，再加入新檔（去重）
                existing = st.session_state.use_data_name if isinstance(st.session_state.use_data_name, list) else []
                new_list = list(dict.fromkeys(existing + saved_names))
                st.cache_data.clear()
                paths = ([DEFAULT_FILE] if st.session_state.include_default else []) + [os.path.join(DATA_FOLDER, f) for f in new_list if os.path.exists(os.path.join(DATA_FOLDER, f))]
                st.session_state.current_data = read_excel_list(paths)
                st.session_state.use_data_name = new_list if new_list else ["DEFAULT"]
                st.session_state.include_default = st.session_state.include_default if new_list else True
                save_data_state("active" if new_list else "default", new_list if new_list else ["FAQ_Default.xlsx"])
                st.success(f"✅ 已加入 {len(saved_names)} 個檔案")
                st.rerun()

        # 使用者手動點擊「X」移除檔案時的重置
        if not uploaded_files:
            # 清空上傳控件的已處理名單，允許再次上傳同名檔案
            st.session_state.processed_files = []
        st.session_state.include_default = st.checkbox("使用預設資料庫", value=st.session_state.include_default, help="是否包含預設資料庫")

        available_files = list_available_files()
        selected = []
        for fn in available_files:
            checked = st.checkbox(fn, value=(fn in (st.session_state.use_data_name or [])), key=f"chk_{fn}")
            if checked:
                selected.append(fn)

        # 顯示目前生效檔案
        try:
            names_list = (selected or [])
            if st.session_state.include_default:
                names_list = ["DEFAULT"] + names_list
            names_str = ", ".join(names_list)
        except Exception:
            names_str = "DEFAULT" if st.session_state.include_default else ""
        st.caption(f"目前生效檔案：{names_str}")

        st.divider()
        # 檔案刪除區（合併於此摺疊）
        delete_candidates = st.multiselect("選擇要刪除的檔案", options=available_files, key="del_candidates")
        if st.button("刪除選擇檔案", key="btn_delete_files", help="注意：刪除後無法復原！", type="primary"):
            deleted, failed = [], []
            for fn in delete_candidates:
                path = os.path.join(DATA_FOLDER, fn)
                try:
                    if os.path.isfile(path):
                        os.remove(path)
                        deleted.append(fn)
                    else:
                        failed.append(fn)
                except Exception:
                    failed.append(fn)

            if deleted:
                # 從使用名單與已處理名單移除
                use_list = st.session_state.use_data_name if isinstance(st.session_state.use_data_name, list) else []
                st.session_state.use_data_name = [f for f in use_list if f not in deleted]
                processed = st.session_state.processed_files if isinstance(st.session_state.processed_files, list) else []
                st.session_state.processed_files = [f for f in processed if f not in deleted]

                # 重新載入資料
                st.cache_data.clear()
                load_paths = ([DEFAULT_FILE] if st.session_state.include_default else []) + [os.path.join(DATA_FOLDER, f) for f in st.session_state.use_data_name if os.path.exists(os.path.join(DATA_FOLDER, f))]
                if load_paths:
                    st.session_state.current_data = read_excel_list(load_paths)
                    save_data_state("default" if (st.session_state.include_default and not st.session_state.use_data_name) else "active", st.session_state.use_data_name if st.session_state.use_data_name else ["FAQ_Default.xlsx"])
                else:
                    st.session_state.include_default = True
                    st.session_state.current_data = read_excel_sheets(DEFAULT_FILE)
                    save_data_state("default", ["FAQ_Default.xlsx"])

                st.success(f"🗑️ 已刪除 {len(deleted)} 個檔案")
                st.rerun()

            if failed:
                st.warning(f"無法刪除：{', '.join(failed)}")

    # 若選擇與現狀不同，更新資料與狀態檔
    if set(selected) != set(st.session_state.use_data_name or [] ) or st.session_state.current_data is None:
        st.session_state.use_data_name = selected
        st.cache_data.clear()
        load_paths = ([DEFAULT_FILE] if st.session_state.include_default else []) + [os.path.join(DATA_FOLDER, f) for f in selected]
        if load_paths:
            st.session_state.current_data = read_excel_list(load_paths)
            save_data_state("default" if (st.session_state.include_default and not selected) else "active", selected)
        else:
            # 無選擇時載入預設
            st.session_state.include_default = True
            st.session_state.current_data = read_excel_sheets(DEFAULT_FILE)
            save_data_state("default", ["FAQ_Default.xlsx"])

    # （已合併）

    # （移至摺疊區塊內）

# --- 5. 生成 System Prompt ---
# 確保 context_data 永遠對應到目前選用的資料 (current_data)
context_text = format_data_for_ai(st.session_state.current_data)
system_prompt = f"""
<角色>你是一名客服人員的專屬助理，可協助客服人員查詢客戶資訊與相關資料並生成建議的回覆的服務</角色>
<任務>
    1. 請先辨識提問，是來自客服人員的提問，還是屬於用戶的提問：
        - 若是來自客服人員的提問，請根據提問內容決定是查詢客戶資訊，還是查詢相關資料，並根據查詢結果生成建議的回覆。
        - 若是來自客戶的提問，請根據<對話歷史>中是否有查詢過客戶資訊，若有，則將客戶姓名帶入回覆中；若沒有，則統一稱為「使用者」，並根據提問內容生成建議的回覆。
    2. 若是查詢客戶資訊，請呼叫 MCP 工具 get_base_info(username) 取得；若在查詢後有生成建議回覆，則將客戶資訊帶入整合；若沒有，直接顯示查詢到的客戶資訊即可。
    3. 若是查詢相關資料，請先使用<資料>（知識庫）內容。若<對話歷史>與<資料>不足以回答，才呼叫工具 get_chat_history_tool 以目前提問抽取的 2–4 個關鍵詞檢索「其他對話」，限制回傳筆數（建議 ≤10），並僅摘錄必要重點。
</任務>
<資料來源與分層>
    1. 優先依據<對話歷史>與<資料>回答。
    2. 其他對話僅作為參考，不代表本次使用者；若與<對話歷史>矛盾，一律以<對話歷史>為準，並簡要說明矛盾點。
    3. 若引用其他對話，請在回覆中加入「其他對話參考」區塊，逐點列出：[file: 檔名] + 80–200 字摘要 + 與本題的關聯理由。
</資料來源與分層>
<限制>
    1. 生成建議的回覆時，需使用``` 區塊必須完整開始並完整結束，區塊結束後，後續說明文字請以一般純文字輸出。
    2. 生成建議的回覆時，請只使用中文文字及數字，不得使用粗體、斜體、底線等格式。
    3. 生成建議的回覆時，清楚、耐心、循序地回應使用者提問，除非使用者明確要求，否則請避免：
        - 長篇說明
        - 顯示程式碼
        - 使用專業縮寫、用語
        - 解釋系統運作原理或展示技術細節
    4. 每次生成建議的回覆時請依照以下流程：
        - 以"OOO您好:" 開頭，若<對話歷史>中有查詢客戶資訊則將客戶姓名帶入，若沒有則統一稱為使用者。
        - 簡要重述使用者問題進行確認，若提問資訊過少，資料中亦無類似的問題，則可引導使用者提供更多資訊。
        - 根據提問提供具體的處理步驟、原因說明或後續行動。
        - 以簡短的關心或確認作為結尾。
    5. 引用其他對話時，不得等同於本次使用者；不得大量貼上原文，需做摘要；每點≤200字。
    6. 除非本次僅為「查詢客戶資訊」並依<查詢客戶資訊回應格式>作答，否則回覆必須在「生成建議回覆回應格式」中列出至少一則參考資料；參考資料來源優先使用<資料>（知識庫），其次可引用「其他對話參考」中的內容，且每則需標示文件名稱與摘錄內容。
</限制>
<生成建議回覆回應格式>
    其他對話參考（如有）:
    - [file: 檔名] 摘要與關聯理由
    - ...

    ---

    - 參考資料1
        - {{參考資料文件名稱}}
        - {{參考資料文件內容}}

    ---

    - 參考資料2
        - {{參考資料文件名稱}}
        - {{參考資料文件內容}}

    ---

    建議回應:

    ```
    {{建議的回應}}
    ```
</生成建議回覆回應格式>
<查詢客戶資訊回應格式>
    - 客戶姓名
    - 裝置世代: {{RouteB or 非RouteB}}
    - 社區: {{社區}}
    - 地區: {{行政區}}
    - 持有電器: {{所持有電器}}
</查詢客戶資訊回應格式>
<資料>{context_text}</資料>
"""

# --- 6. 主介面顯示 ---
st.title("Customer Service Wingman")
st.caption("Version: v2.2.2")

# 顯示現有的對話紀錄
for message in st.session_state.messages:
    avatar_icon = USER_AVATAR if message["role"] == "user" else CSW_AVATAR
    with st.chat_message(message["role"], avatar=avatar_icon):
        st.markdown(message["content"])

# --- 7. 對話邏輯 ---
if prompt := st.chat_input("請問我有什麼可以協助的嗎?"):

    # 檢查驗證
    if not api_valid:
        st.error("驗證失敗：請檢查後在左側選單重新輸入 API Key")
        st.stop()
    if not st.session_state.current_data:
        st.error("缺少資料庫資料")
        st.stop()

    # 顯示使用者訊息
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user", avatar=USER_AVATAR):
        st.markdown(prompt)

    # 呼叫 Akasha 回覆
    with st.chat_message("assistant", avatar=CSW_AVATAR):
        with st.spinner("思考中..."):
            try:
                # 使用與目前執行的 Python 同一個解譯器
                python_cmd = sys.executable or "python"
                # 以目前檔案位置為基準定位 tools 目錄
                project_root = os.path.dirname(os.path.abspath(__file__))
                tools_path = os.path.join(project_root, "tools")
                if not os.path.exists(tools_path):
                    st.error(f"找不到工具腳本：{tools_path}")
                    st.stop()
                connection_info = {
                    "get_user_info_tool": {
                        "command": python_cmd,
                        "args": [os.path.join(tools_path, "get_user_info.py")],
                        "transport": "stdio",
                    },
                    "get_chat_history_tool": {
                        "command": python_cmd,
                        "args": [os.path.join(tools_path, "get_chat_history.py")],
                        "transport": "stdio",
                    }
                }
                
                agent = akasha.agents(
                    model=config["model_name"],
                    temperature=0.7,
                    max_input_tokens=20000,
                    max_output_tokens=20000,
                    verbose=True
                )
                history_text = get_history_string(st.session_state.history_list)
                final_prompt = (
                    system_prompt +
                    f"\n<提問>\n{prompt}\n</提問>" + 
                    f"\n<對話歷史>\n{history_text}\n</對話歷史>"
                )
                response = agent.mcp_agent(connection_info, final_prompt)
                resp_out = normalize_response_text(response)
                st.markdown(resp_out)

                # --- Token 管理與修剪 --- 
                st.session_state.history_list.append({"q": prompt, "a": resp_out})
                
                # 更新並計算 Token
                current_h_text = get_history_string(st.session_state.history_list)
                total_content = system_prompt + prompt + current_h_text
                
                # 迴圈修剪
                while ah.myTokenizer.compute_tokens(total_content, config["model_name"]) > 20000 and len(st.session_state.history_list) > 1:
                    st.session_state.history_list.pop(0)
                    current_h_text = get_history_string(st.session_state.history_list)
                    total_content = system_prompt + prompt + current_h_text

                # 存回 messages 用於顯示
                st.session_state.messages.append({"role": "assistant", "content": resp_out})
                # 自動儲存對話紀錄（首次訊息時建立檔案）
                prev_active = get_chat_active_file()
                new_path = save_chat_log(create_if_missing=True)
                # 若為第一次建立新對話檔案，重新整理以刷新側邊欄列表與預設選取
                if not prev_active and new_path:
                    st.rerun()
            except Exception as err:
                err_type = err.__class__.__name__ # 取得錯誤的class 名稱
                info = err.args[0] # 取得詳細內容
                detains = traceback.format_exc() # 取得完整的tracestack
                n1, n2, n3 = sys.exc_info() #取得Call Stack
                lastCallStack =  traceback.extract_tb(n3)[-1] # 取得Call Stack 最近一筆的內容
                fn = lastCallStack [0] # 取得發生事件的檔名
                lineNum = lastCallStack[1] # 取得發生事件的行數
                funcName = lastCallStack[2] # 取得發生事件的函數名稱
                errMesg = f"FileName: {fn}, lineNum: {lineNum}, Fun: {funcName}, reason: {info}, trace:\n {traceback.format_exc()}"
                print(errMesg)
                st.error(f"模型呼叫失敗: {str(err)}")